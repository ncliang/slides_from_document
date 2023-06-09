import argparse
import os

from langchain import PromptTemplate, LLMChain
from langchain.chat_models import ChatOpenAI
from langchain.document_loaders import UnstructuredURLLoader, UnstructuredFileLoader
from langchain.schema import Document
from pptx import Presentation

from markdown_slides_output_parser import MarkdownSlidesOutputParser

PROMPT_TEMPLATE = """用以下內容產生一份沒有圖片的投影片。{format_instructions}

內容：

{content}
"""


def load_documents(cmdline_args):
    if cmdline_args.url:
        return UnstructuredURLLoader(urls=[cmdline_args.url]).load()
    elif cmdline_args.file:
        return UnstructuredFileLoader(cmdline_args.file).load()
    elif cmdline_args.content:
        return [Document(page_content=cmdline_args.content)]
    else:
        raise ValueError(f"Unable to create loader from args: {cmdline_args}")


parser = argparse.ArgumentParser(description="Generate slides from provided content using ChatGPT")
parser.add_argument("--output", type=str, required=True, help="Location of generated slide")
parser.add_argument("--template", type=str, help="Optional pptx template to use for slide generation")

group = parser.add_mutually_exclusive_group(required=True)
group.add_argument("--content", type=str, help="Content to use for slide generation")
group.add_argument("--url", type=str, help="URL to content to use for slide generation")
group.add_argument("--file", type=str, help="A file to use as content for slide generation")

args = parser.parse_args()

documents = load_documents(args)

print(f"Loaded content with {len(documents[0].page_content)} characters")

llm = ChatOpenAI(temperature=0)
output_parser = MarkdownSlidesOutputParser()
prompt = PromptTemplate(
    input_variables=["content"],
    partial_variables={"format_instructions": output_parser.get_format_instructions()},
    template=PROMPT_TEMPLATE,
)

print(f"Sending content to ChatGPT for slide generation.")
chain = LLMChain(llm=llm, prompt=prompt)
result = chain.run(documents[0].page_content)
print(f"Response from ChatGPT: \n{result}")
slides = output_parser.parse(result)
assert slides

print(f"Generated presentation with {len(slides)} slides. Title slide is `{slides[0].title}`")
output_file = os.path.join(args.output, "out.pptx")

prs = Presentation()
if args.template:
    print(f"Reading from specified template file {args.template}")
    prs = Presentation(args.template)

for slide in slides:
    slide.add_slide(prs)

print(f"Writing presentation to {output_file}")
prs.save(output_file)
