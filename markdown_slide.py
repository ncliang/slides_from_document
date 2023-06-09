import json
from typing import List

from pptx.presentation import Presentation

from util import get_content_only


class SlideWithTitle(object):
    title: str

    def __init__(self):
        self.title = ""

    def add_slide(self, presentation: Presentation):
        raise NotImplementedError()

    def __repr__(self):
        return json.dumps(self, default=lambda o: o.__dict__, sort_keys=True, indent=2)


class SlideWithSubtitle(SlideWithTitle):
    subtitle: str

    def __init__(self):
        super().__init__()
        self.subtitle = ""

    def add_slide(self, presentation: Presentation):
        title_slide_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = self.title
        subtitle.text = self.subtitle


class BulletPoint(object):
    text: str
    level: int  # the indentation level

    def __init__(self, text, level):
        self.text = text
        self.level = level

    @staticmethod
    def from_bullet_point_line(line):
        leading_spaces = len(line) - len(line.lstrip(" "))
        return BulletPoint(
            get_content_only(line.lstrip(" ")),
            int(leading_spaces / 2)
        )

    def __repr__(self):
        return json.dumps(self, default=lambda o: o.__dict__, sort_keys=True, indent=2)


class SlideWithBulletPoints(SlideWithTitle):
    bullet_points: List[BulletPoint]

    def __init__(self):
        super().__init__()
        self.bullet_points = []

    def add_slide(self, presentation: Presentation):
        bullet_slide_layout = presentation.slide_layouts[1]

        slide = presentation.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = self.title

        tf = body_shape.text_frame
        tf.text = self.bullet_points.pop(0).text
        for bp in self.bullet_points:
            p = tf.add_paragraph()
            p.text = bp.text
            p.level = int(bp.level)
